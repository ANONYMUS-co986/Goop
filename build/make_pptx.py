#!/usr/bin/env python3
"""Pack build/slides/slide-NN.png into Project_Verde_Presentation.pptx.

Each slide is a 300-DPI full-bleed raster captured from the same Chromium
pipeline that renders the FX PDF, so the deck is a pixel-clone of
Project_Verde_Documentation.pdf.

Size discipline: text rasters stay lossless PNG; big photographic rasters
(> PNG_KEEP bytes) become JPEG q93 — visually indistinguishable at 300 DPI,
but keeps the deck upload-friendly.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

BUILD = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BUILD)
SLIDES = os.path.join(BUILD, "slides")
OUT = os.path.join(ROOT, "Project_Verde_Presentation.pptx")

TARGET_W, TARGET_H = 2480, 3508   # A4 portrait @ 300 DPI
PNG_KEEP = 1_500_000              # bytes; above -> JPEG q93

files = sorted(f for f in os.listdir(SLIDES) if f.endswith(".png"))
assert len(files) == 33, f"expected 33 slide rasters, found {len(files)}"

prepared, kept_png, made_jpg = [], 0, 0
for f in files:
    p = os.path.join(SLIDES, f)
    im = Image.open(p).convert("RGB")
    resized = im.size != (TARGET_W, TARGET_H)
    if resized:
        im = im.resize((TARGET_W, TARGET_H), Image.LANCZOS)
    if os.path.getsize(p) > PNG_KEEP:
        jp = p[:-4] + ".jpg"
        im.save(jp, "JPEG", quality=93, optimize=True, subsampling=1)
        prepared.append(jp)
        made_jpg += 1
    else:
        if resized:
            im.save(p, "PNG", optimize=True)
        prepared.append(p)
        kept_png += 1

prs = Presentation()
prs.slide_width = Emu(210 * 36000)    # 210 mm
prs.slide_height = Emu(297 * 36000)   # 297 mm
blank = prs.slide_layouts[6]
for img in prepared:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)

cp = prs.core_properties
cp.title = "Project Verde — Smart IoT Irrigation & Plant-Care System"
cp.author = "Aarav Choudhary & Anuj (Class X) — DAV ACON 5, 2026"
cp.subject = "DAV ACON 5 Tech Exhibition 2026 — competition deck"
cp.keywords = "Project Verde, IoT, ESP32, Firebase, irrigation, plant care, DAV ACON 5"
cp.comments = "33-slide exhibition deck. Pixel-clone of the Project Verde FX documentation."
cp.category = "Exhibition Submission"
cp.last_modified_by = "Project Verde build pipeline"
prs.save(OUT)

# ---- read-back verification -----------------------------------------------
chk = Presentation(OUT)
assert len(chk.slides) == 33, f"slide count mismatch: {len(chk.slides)}"
for i, s in enumerate(chk.slides, 1):
    pics = [sh for sh in s.shapes if sh.shape_type == 13]  # PICTURE
    assert len(pics) == 1, f"slide {i}: expected 1 picture, got {len(pics)}"
    pic = pics[0]
    assert pic.width == chk.slide_width and pic.height == chk.slide_height, \
        f"slide {i}: picture does not fill the slide"
    assert pic.left == 0 and pic.top == 0, f"slide {i}: picture offset"
sz = os.path.getsize(OUT)
print(f"OK  {OUT}")
print(f"    slides: {len(chk.slides)}  |  {kept_png} lossless PNG + {made_jpg} JPEG q93")
print(f"    slide size: {chk.slide_width} x {chk.slide_height} EMU (210x297 mm A4 portrait)")
print(f"    file size: {sz/1_048_576:.1f} MB")
